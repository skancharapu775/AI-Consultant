"""Generate PowerPoint deck from analysis results."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import io
from typing import List, Dict, Any
from datetime import datetime


def generate_deck(
    pnl_data: List[Dict[str, Any]],
    diagnostics: Dict[str, Any],
    initiatives: List[Dict[str, Any]],
    data_completeness: Dict[str, Any] = None,
) -> bytes:
    """Generate PowerPoint deck and return as bytes."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Executive Summary
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_executive_summary_slide(slide1, initiatives)

    # Slide 2: P&L + EBITDA Trend
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_pnl_trend_slide(slide2, pnl_data)

    # Slide 3: Cost Structure
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_cost_structure_slide(slide3, pnl_data)

    # Slide 4: Top Initiative Deep Dive
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    if initiatives:
        top_initiative = sorted(initiatives, key=lambda x: x.get("rank", 999))[0]
        add_initiative_detail_slide(slide4, top_initiative)

    # Slide 5: Roadmap or Data Gaps
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    if data_completeness and data_completeness.get("data_gaps"):
        add_data_gaps_slide(slide5, data_completeness)
    else:
        add_roadmap_slide(slide5, initiatives)

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def add_title_and_subtitle(slide, title_text: str, subtitle_text: str = ""):
    """Add title and subtitle to a slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True

    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


def add_executive_summary_slide(slide, initiatives: List[Dict[str, Any]]):
    """Add executive summary slide."""
    add_title_and_subtitle(slide, "Executive Summary", "Top Improvement Initiatives")

    top_5 = sorted(initiatives, key=lambda x: x.get("rank", 999))[:5]
    y_start = 1.8

    for idx, init in enumerate(top_5):
        y_pos = y_start + (idx * 1.0)
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.5), Inches(0.9))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        impact_str = f"${init.get('impact_low', 0)/1000:.0f}K - ${init.get('impact_high', 0)/1000:.0f}K"
        p = text_frame.paragraphs[0]
        p.text = f"{init.get('rank', 'N/A')}. {init.get('title', 'N/A')} — {impact_str} annual impact"
        p.font.size = Pt(14)
        p.font.bold = True


def add_pnl_trend_slide(slide, pnl_data: List[Dict[str, Any]]):
    """Add P&L trend chart slide."""
    add_title_and_subtitle(slide, "P&L & EBITDA Trend")

    if not pnl_data:
        return

    # Create chart
    months = [d["month"] for d in pnl_data]
    revenue = [d["revenue"] for d in pnl_data]
    ebitda = [d["ebitda"] for d in pnl_data]

    fig, ax1 = plt.subplots(figsize=(8, 4))
    ax1.set_xlabel("Month")
    ax1.set_ylabel("Revenue ($)", color="blue")
    ax1.plot(months, revenue, color="blue", marker="o", label="Revenue")
    ax1.tick_params(axis="y", labelcolor="blue")

    ax2 = ax1.twinx()
    ax2.set_ylabel("EBITDA ($)", color="green")
    ax2.plot(months, ebitda, color="green", marker="s", label="EBITDA")
    ax2.tick_params(axis="y", labelcolor="green")

    plt.title("Revenue and EBITDA Trend")
    plt.xticks(rotation=45)
    plt.tight_layout()

    # Save to bytes
    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format="png", dpi=150)
    img_buffer.seek(0)
    plt.close()

    # Add image to slide
    slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.8), width=Inches(8))


def add_cost_structure_slide(slide, pnl_data: List[Dict[str, Any]]):
    """Add cost structure stacked area chart."""
    add_title_and_subtitle(slide, "Operating Expense Structure")

    if not pnl_data:
        return

    months = [d["month"] for d in pnl_data]
    sales_mkt = [d.get("opex_sales_marketing", 0) for d in pnl_data]
    rnd = [d.get("opex_rnd", 0) for d in pnl_data]
    gna = [d.get("opex_gna", 0) for d in pnl_data]
    other = [d.get("opex_other", 0) for d in pnl_data]

    fig, ax = plt.subplots(figsize=(8, 4))
    ax.stackplot(
        months,
        sales_mkt,
        rnd,
        gna,
        other,
        labels=["Sales & Marketing", "R&D", "G&A", "Other"],
        alpha=0.7,
    )
    ax.set_xlabel("Month")
    ax.set_ylabel("Operating Expenses ($)")
    ax.set_title("Operating Expense Structure Over Time")
    ax.legend(loc="upper left")
    plt.xticks(rotation=45)
    plt.tight_layout()

    img_buffer = io.BytesIO()
    plt.savefig(img_buffer, format="png", dpi=150)
    img_buffer.seek(0)
    plt.close()

    slide.shapes.add_picture(img_buffer, Inches(1), Inches(1.8), width=Inches(8))


def add_initiative_detail_slide(slide, initiative: Dict[str, Any]):
    """Add detailed view of top initiative."""
    add_title_and_subtitle(slide, f"Top Initiative: {initiative.get('title', 'N/A')}")

    y_pos = 1.8
    details = [
        f"Category: {initiative.get('category', 'N/A')}",
        f"Owner: {initiative.get('owner', 'N/A')}",
        f"Impact: ${initiative.get('impact_low', 0):,.0f} - ${initiative.get('impact_high', 0):,.0f} annually",
        f"Time to Value: {initiative.get('time_to_value_weeks', 0)} weeks",
        f"Implementation Cost: ${initiative.get('implementation_cost_estimate', 0):,.0f}",
        f"Risk Level: {initiative.get('risk_level', 'N/A')}",
        f"Confidence: {initiative.get('confidence', 0)*100:.0f}%",
    ]

    description = initiative.get("description", "")
    if description:
        details.append(f"\nDescription:\n{description}")

    text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = "\n".join(details)

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(12)


def add_roadmap_slide(slide, initiatives: List[Dict[str, Any]]):
    """Add roadmap/sequencing slide."""
    add_title_and_subtitle(slide, "Implementation Roadmap", "Suggested Sequencing")

    top_5 = sorted(initiatives, key=lambda x: x.get("rank", 999))[:5]
    y_start = 1.8

    for idx, init in enumerate(top_5):
        y_pos = y_start + (idx * 0.9)
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.5), Inches(0.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        time_str = f"{init.get('time_to_value_weeks', 0)} weeks"
        p = text_frame.paragraphs[0]
        p.text = (
            f"Q{idx+1}: {init.get('title', 'N/A')} "
            f"(Impact: ${init.get('impact_low', 0)/1000:.0f}K-${init.get('impact_high', 0)/1000:.0f}K, "
            f"Time: {time_str})"
        )
        p.font.size = Pt(13)


def add_data_gaps_slide(slide, data_completeness: Dict[str, Any]):
    """Add data gaps slide."""
    add_title_and_subtitle(slide, "Data Gaps / What Would Improve Confidence")

    y_pos = 1.8
    text_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.5), Inches(5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True

    data_gaps = data_completeness.get("data_gaps", [])
    text_lines = ["Missing or Incomplete Data:", ""]
    
    for gap in data_gaps[:8]:  # Limit to 8 items
        text_lines.append(f"• {gap}")
    
    if len(data_gaps) > 8:
        text_lines.append(f"... and {len(data_gaps) - 8} more")
    
    text_lines.append("")
    text_lines.append("Impact on Analysis:")
    text_lines.append("• Missing optional datasets reduce confidence in initiative sizing")
    text_lines.append("• Some initiative types may be disabled or have wider impact ranges")
    text_lines.append(f"• Completeness score: {data_completeness.get('completeness_score', 0)*100:.0f}%")

    text_frame.text = "\n".join(text_lines)
    
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(12)
